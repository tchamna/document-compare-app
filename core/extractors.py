"""File-type extractors: PPTX, DOCX (with corrupt-file fallback), TXT."""

from __future__ import annotations

import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from docx import Document

from core.helpers import normalize_whitespace, split_into_lines, is_digits_only


# ------------------------------------------------------------------ #
#  Internal PPTX helpers
# ------------------------------------------------------------------ #
def _shape_text_lines(shape) -> List[str]:
    """Extract text from a single PPTX shape."""
    if not getattr(shape, "has_text_frame", False):
        return []
    raw = shape.text_frame.text or ""
    if not raw.strip():
        return []
    return split_into_lines(raw)


# ------------------------------------------------------------------ #
#  Public extractors
# ------------------------------------------------------------------ #
def extract_slide_lines(pptx_path: Path) -> Dict[int, List[str]]:
    """Extract text per slide from a PPTX file.

    Returns ``{slide_number: [line, ...]}``.
    """
    prs = Presentation(str(pptx_path))
    slide_map: Dict[int, List[str]] = {}

    for i, slide in enumerate(prs.slides, start=1):
        lines: List[str] = []
        for shape in slide.shapes:
            lines.extend(_shape_text_lines(shape))
        lines = [ln for ln in lines if ln.strip() and not is_digits_only(ln)]
        slide_map[i] = lines

    return slide_map


def extract_docx_lines(docx_path: Path) -> List[str]:
    """Extract paragraphs from a DOCX file.

    Falls back to raw XML extraction when the file contains corrupt
    embedded media (``BadZipFile``).
    """
    try:
        doc = Document(str(docx_path))
        paragraphs = [para.text for para in doc.paragraphs]
    except zipfile.BadZipFile:
        print(
            "\u26a0\ufe0f  Fichier DOCX corrompu (image invalide). "
            "Extraction du texte brut..."
        )
        ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
        paragraphs = []
        with zipfile.ZipFile(str(docx_path), "r") as zf:
            with zf.open("word/document.xml") as f:
                tree = ET.parse(f)
        for p_elem in tree.iter(f"{ns}p"):
            texts = [t.text for t in p_elem.iter(f"{ns}t") if t.text]
            paragraphs.append("".join(texts))
    except PermissionError:
        raise PermissionError(
            f"Le fichier est verrouillé (probablement ouvert dans Word) :\n"
            f"  {docx_path}\nFermez-le et réessayez."
        )
    except Exception as e:
        raise RuntimeError(
            f"Impossible d'ouvrir le fichier DOCX : {docx_path}\n  Erreur : {e}"
        )

    all_lines: List[str] = []
    for text in paragraphs:
        normalized = normalize_whitespace(text)
        if normalized and not is_digits_only(normalized):
            all_lines.append(normalized)
    return all_lines


def extract_txt_lines(txt_path: Path) -> List[str]:
    """Extract lines from a plain-text file."""
    with open(txt_path, "r", encoding="utf-8") as f:
        content = f.read()

    all_lines: List[str] = []
    for block in content.split("\n\n"):
        lines = split_into_lines(block)
        lines = [ln for ln in lines if ln.strip() and not is_digits_only(ln)]
        all_lines.extend(lines)
    return all_lines


def extract_text_lines(file_path: Path) -> List[str]:
    """Auto-detect file type and return a flat list of lines.

    For PPTX files, use :func:`extract_slide_lines` directly.
    """
    suffix = file_path.suffix.lower()
    if suffix == ".docx":
        return extract_docx_lines(file_path)
    if suffix == ".txt":
        return extract_txt_lines(file_path)
    if suffix == ".doc":
        raise ValueError(
            ".doc files (legacy Word) not yet supported. "
            "Please convert to .docx or use .txt"
        )
    raise ValueError(
        f"Unsupported file type: {suffix}. Supported: .pptx, .docx, .txt"
    )
